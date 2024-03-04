import os
from docx import Document
from PIL import Image
from io import BytesIO
from pptx import Presentation

def extract_images_from_docx(docx_path, output_images_dir):
    doc = Document(docx_path)
    # Create the output images directory if it doesn't exist
    if not os.path.exists(output_images_dir):
        os.makedirs(output_images_dir)
    image_index = 1
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_data = rel.target_part.blob
            image = Image.open(BytesIO(image_data))
            image.save(os.path.join(output_images_dir, f'image_{image_index}.png'))
            image_index += 1

def extract_images_from_pptx(pptx_path, output_images_dir):
    prs = Presentation(pptx_path)
    # Create the output images directory if it doesn't exist
    if not os.path.exists(output_images_dir):
        os.makedirs(output_images_dir)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_image_frame:
                image = shape.image
                image_bytes = image.blob
                image_data = BytesIO(image_bytes)
                img = Image.open(image_data)
                img.save(os.path.join(output_images_dir, f'image_{image_index}.png'))
                image_index += 1

