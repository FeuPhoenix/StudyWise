import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
from PIL import Image
from io import BytesIO
from summarizer import Summarizer
api_key = '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'

#from flashcards_from_pdf.flashcard_creator import runFlashcards
def getFileNameFromPathWithOutExtension(input_string):
    # Find the last occurrence of '/'
    last_slash_index = input_string.rfind('/')
    
    # Slice the string from the character after the last '/'
    # If '/' is not found, rfind returns -1, and slicing starts from index 0
    result_string = input_string[last_slash_index + 1:]
    result_string=result_string.replace('.mp4','')
    result_string=result_string.replace('.docx','')
    result_string=result_string.replace('.doc','')
    result_string=result_string.replace('.pptx','')
    result_string=result_string.replace('.ppt','')
    result_string=result_string.replace('.pdf','')
    return result_string
def delete_page_breaks(file_path):
    # Read the content of the file
    with open(file_path, 'r',encoding='utf-8') as file:
        content = file.read()

    # Remove "--- Page Break ---" occurrences
    cleaned_content = content.replace("--- Page Break ---", "")

    # Write the cleaned content back to the file
    with open(file_path, 'w') as file:
        file.write(cleaned_content)
def clean_text(text):
    
    allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
    
    cleaned_text = re.sub(allowed_pattern, '', text)    
    return cleaned_text
def read_text_file(file_path):
    with open(file_path, 'r',encoding='utf-8') as file:
        return file.read()

def extract_text_from_word(file_path):
 
    # Load the Word document
    doc = Document(file_path)
    
    # Initialize an empty string to store the extracted text
    text = ''
    
    # Iterate through paragraphs in the document and concatenate the text
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    text=clean_text(text)
    return text
def extract_text_from_pptx(prs):
    text = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'

    return text

def extract_images_from_pptx(prs, output_images_dir):
    # Create the output images directory if it doesn't exist
    if not os.path.exists(output_images_dir):
        os.makedirs(output_images_dir)

    image_index = 1
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_data = BytesIO(image_bytes)
                img = Image.open(image_data)
                img.save(os.path.join(output_images_dir, f'image_{slide_number}_{shape_number}.png'))
                image_index += 1

def save_text_to_txt(text, txt_path):
    # Save the extracted text to a text file
    with open(txt_path, 'w', encoding='utf-8') as file:
        file.write(text)


def split_text(text, max_chunk_size=3800):  # Reduced max_chunk_size for safety
    
    sentences = re.split(r'(?<=[.!?]) +', text)
    current_chunk = ""
    for sentence in sentences:
        # Check if adding the next sentence would exceed the max chunk size EZ
        if len(current_chunk) + len(sentence) > max_chunk_size:
            yield current_chunk
            current_chunk = sentence  # Start a new chunk EZ
        else:
            current_chunk += " " + sentence
    yield current_chunk  # Yield the last chunk

def get_Long_summary_from_gpt3(text, api_key):
    openai.api_key = api_key
    summaries = []

    for chunk in split_text(text):
        while True:  # Keep trying until successful
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant."},
                        {"role": "user", "content": f"Summarize the following text into concise bullet points, focusing only on essential information: \n{chunk}"}
                    ]
                )
                summaries.append(response['choices'][0]['message']['content'].strip())
                break  # Exit the loop if successful
            except openai.error.RateLimitError as e:
                print("Rate limit exceeded, waiting to retry...")
                time.sleep(20)  # Wait for 20 seconds before retrying

    full_summary = ' '.join(summaries)
    return full_summary



def create_json_with_Long_summary(json_file_path, summary):
    # Extracting the first sentence for the title
    title_match = re.match(r"([^.]*.)", summary)
    title = title_match.group(0).strip() if title_match else "Summary"

    # Creating a dictionary for JSON data
    summary_data = {
        
        'summary': summary
    }

    # Writing the dictionary to a JSON file
    with open(json_file_path, 'w') as file:
        json.dump(summary_data, file, indent=4)

def extract_text_from_pdf_plumber(pdf_path, txt_file_path):
    # Open the PDF file
    if not os.path.exists(txt_file_path):
     print("does not exists")
        # Create the directory, including any necessary parent directories
     with pdfplumber.open(pdf_path) as pdf:
        # Open the text file in write mode
        with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
            # Iterate through each page in the PDF
            for page in pdf.pages:
                # Extract text from the page
                text = page.extract_text()
                
                # Write the text to the text file, if text was found
                if text:
                    txt_file.write(text)
                    
                    #  add a page break in the text file
            print(f"Text extracted and saved to {txt_file_path}")     
def is_Document(file):
    
    
    if os.path.isfile(file) and file.endswith('.pdf'):
        PDFFile=file
        filename=getFileNameFromPathWithOutExtension(file)
        text_file_path = f'Summaries/transcribed_text_From_{filename}.txt'
        #json_file_path = f'Summaries/summary_From_PDF{filename}.json'

        extract_text_from_pdf_plumber(PDFFile,text_file_path)
        text = read_text_file(text_file_path)
        summary = get_Long_summary_from_gpt3(text, api_key)

        # create_json_with_Long_summary(json_file_path, summary)
        

        text=read_text_file(text_file_path)
    
        text=clean_text(text)

        
        with open(text_file_path, 'w') as file:
                file.write(summary)
                delete_page_breaks(text_file_path)
        file_data=read_text_file(text_file_path)
        summary_data = {
                'long_summary': file_data
            }

        with open(f"assets/output_files/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 

    elif file.endswith('.ppt') or file.endswith('.pptx'):
        pptx_path = file
        filename=getFileNameFromPathWithOutExtension(file)
        #output_text_path = f'assets/output_files/{filename}.txt'

        

        # Load the PowerPoint presentation
        prs = Presentation(pptx_path)

        
        text = extract_text_from_pptx(prs)

        
        
        # runFlashcards(output_text_path)
        summary = get_Long_summary_from_gpt3(text, api_key)
        #create_json_with_Long_summary(json_file_path, summary)
        summary_data = {
                'long_summary': summary
            }

        with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
    elif file.endswith('.doc') or file.endswith('.docx'):
        docx_path = file
        filename=getFileNameFromPathWithOutExtension(file)
        #output_text_path = f'assets/output_files/{filename}.txt'

        #output_images_dir = f'assets/output_files/{filename}_images'
        text=extract_text_from_word(docx_path)
        #runFlashcards(output_text_path)
        summary = get_Long_summary_from_gpt3(text, api_key)
        #create_json_with_Long_summary(json_file_path, summary)
        summary_data = {
                'long_summary': summary
            }

        with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
       

        
       
    
#testing
def main():
    file = "assets/input_files/text-based/test2.pdf"  
    is_Document(file)
#Dido
# def main(file):
#     # file = "assets/input_files/text-based/test2.pptx"  
#     is_Document(file)
if __name__==main():
    main()




