import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
from PIL import Image
from io import BytesIO
from summarizer import Summarizer

def getFileNameFromPathWithOutExtension(input_string):
    last_slash_index = input_string.rfind('/')
    result_string = input_string[last_slash_index + 1:]
    result_string=result_string.replace('.mp4','')
    result_string=result_string.replace('.docx','')
    result_string=result_string.replace('.doc','')
    result_string=result_string.replace('.pptx','')
    result_string=result_string.replace('.ppt','')
    result_string=result_string.replace('.pdf','')
    return result_string

def delete_page_breaks(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    cleaned_content = content.replace("--- Page Break ---", "")
    with open(file_path, 'w') as file:
        file.write(cleaned_content)

def clean_text(text):
    allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
    cleaned_text = re.sub(allowed_pattern, '', text)    
    return cleaned_text

def read_text_file(file_path):
    with open(file_path, 'r',encoding='utf-8') as file:
        return file.read()

def extract_text_from_word(file_path):
    doc = Document(file_path)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    text = clean_text(text)
    return text

def extract_text_from_pptx(prs):
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_images_from_pptx(prs, output_images_dir):
    if not os.path.exists(output_images_dir):
        os.makedirs(output_images_dir)
    image_index = 1
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_data = BytesIO(image_bytes)
                img = Image.open(image_data)
                img.save(os.path.join(output_images_dir, f'image_{slide_number}_{shape_number}.png'))
                image_index += 1

def save_text_to_txt(text, txt_path):
    with open(txt_path, 'w', encoding='utf-8') as file:
        file.write(text)

def split_text(text, max_chunk_size=3800):
    sentences = re.split(r'(?<=[.!?]) +', text)
    current_chunk = ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) > max_chunk_size:
            yield current_chunk
            current_chunk = sentence
        else:
            current_chunk += " " + sentence
    yield current_chunk

def get_Long_summary_from_gpt3(text, api_key):
    openai.api_key = api_key
    summaries = []
    for chunk in split_text(text):
        while True:
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant."},
                        {"role": "user", "content": f"Summarize the following text into concise bullet points, focusing only on essential information: \n{chunk}"}
                    ]
                )
                summaries.append(response['choices'][0]['message']['content'].strip())
                break
            except openai.error.RateLimitError as e:
                print("Rate limit exceeded, waiting to retry...")
                time.sleep(20)
    full_summary = ' '.join(summaries)
    return full_summary

def create_json_with_Long_summary(json_file_path, summary):
    title_match = re.match(r"([^.]*.)", summary)
    title = title_match.group(0).strip() if title_match else "Summary"
    summary_data = {
        'summary': summary
    }
    with open(json_file_path, 'w') as file:
        json.dump(summary_data, file, indent=4)

def extract_text_from_pdf_plumber(pdf_path, txt_file_path):
    if not os.path.exists(txt_file_path):
        with pdfplumber.open(pdf_path) as pdf:
            with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        txt_file.write(text)
            print(f"Text extracted and saved to {txt_file_path}")     

def is_Document(file):
    if os.path.isfile(file) and file.endswith('.pdf'):
        PDFFile=file
        filename=getFileNameFromPathWithOutExtension(file)
        text_file_path = f'Summaries/transcribed_text_From_{filename}.txt'
        json_file_path = f'Summaries/summary_From_PDF{filename}.json'
        api_key = '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'
        extract_text_from_pdf_plumber(PDFFile,text_file_path)
        text = read_text_file(text_file_path)
        model = Summarizer()
        text=clean_text(text)
        result = model(text)
        summary_data = {
            'long_summary': result
        }
        with open(f"assets/output_files/{filename}_summary.json", 'w') as json_file:
            json.dump(summary_data, json_file, indent=4)
            print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
    elif os.path.isfile(file) and (file.endswith('.ppt') or file.endswith('.pptx')):
        pptx_path = file
        filename=getFileNameFromPathWithOutExtension(file)
        output_text_path = f'assets/output_files/{filename}.txt'
        json_file_path = f'Summaries/summary_From_PPT{file}.json'
        prs = Presentation(pptx_path)
        text = extract_text_from_pptx(prs)
        model = Summarizer()
        text=clean_text(text)
        result = model(text)
        summary_data = {
            'long_summary': result
        }
        with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
            json.dump(summary_data, json_file, indent=4)
            print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
    elif os.path.isfile(file) and (file.endswith('.doc') or file.endswith('.docx')):
        docx_path = file
        filename=getFileNameFromPathWithOutExtension(file)
        text=extract_text_from_word(docx_path)
        model = Summarizer()
        result = model(text)
        summary_data = {
            'long_summary': result
        }
        with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
            json.dump(summary_data, json_file, indent=4)
            print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 

<<<<<<< HEAD
def main(file):
    # file = "assets/input_files/text-based/test.pdf"  
    is_Document(file)
if __name__ == "__main__":
    main()
=======
#Testing
def main():
    file = "assets/input_files/text-based/test.pdf"  
    is_Document(file)
if __name__ == "__main__":
    main()
#Dido
# def main(file):
#     #file = "assets/input_files/text-based/test.pdf"  
#     is_Document(file)
>>>>>>> mvc
