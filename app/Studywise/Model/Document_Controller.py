import subprocess
import uuid
from datetime import datetime
from firebase_admin import firestore, credentials, initialize_app
import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
from PIL import Image
from io import BytesIO
from summarizer import Summarizer
from firebase_admin import credentials, storage
from Flash_Cards_Controller import FlashcardsController
from DocumentProcessed_Repo import DocumentProcessed
from Question_Controller import QuestionController
import fitz  # PyMuPDF
import os
import comtypes.client



class DocumentProcessedController:
    def __init__(self,file) :
        self.file=file
        self.Document_Processing(file)
       
    @staticmethod
    # async def create_processed_document(request_data):
    #     """
    #     Creates a new processed document record in Firestore.
    #     :param request_data: The data from the request to create a new processed document.
    #     """
    #     try:
    #         # Extracting necessary information from the request data
    #         processed_material_id = request_data.get('processed_material_id')
    #         material_id = request_data.get('material_id')
    #         generated_summary_file_path = request_data.get('generated_summary_file_path')
    #         generated_text_file_path = request_data.get('generated_text_file_path')
    #         generated_images_file_path = request_data.get('generated_images_file_path')
    #         generated_document_file_path = request_data.get('generated_document_file_path')
            
    #         # Creating a new DocumentProcessed instance
    #         new_document = DocumentProcessed(
    #             processed_material_id=processed_material_id,
    #             material_id=material_id,
    #             generated_summary_file_path=generated_summary_file_path,
    #             generated_text_file_path=generated_text_file_path,
    #             generated_images_file_path=generated_images_file_path,
    #             generated_Document_file_path=generated_document_file_path
    #         )
            
    #         # Adding the new document to Firestore
    #         await new_document.addProcessedMaterialToFirestore()
            
    #         return {"success": True, "message": "Processed document added successfully."}
    #     except Exception as e:
    #         return {"success": False, "message": str(e)}

    
    # Additional methods for updating and deleting processed documents can be added here
    @staticmethod
    def getFileNameFromPathWithOutExtension(input_string):
        last_slash_index = input_string.rfind('/')
        result_string = input_string[last_slash_index + 1:]
        result_string=result_string.replace('.mp4','')
        result_string=result_string.replace('.docx','')
        result_string=result_string.replace('.doc','')
        result_string=result_string.replace('.pptx','')
        result_string=result_string.replace('.ppt','')
        result_string=result_string.replace('.pdf','')
        return result_string

    @staticmethod
    def clean_text(text):
        allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
        cleaned_text = re.sub(allowed_pattern, '', text)    
        return cleaned_text

    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r',encoding='utf-8', errors='ignore') as file:
            return file.read()

    @staticmethod
    def extract_text_from_word(file_path):
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        text = DocumentProcessedController.clean_text(text)
        return text

    @staticmethod
    def extract_text_from_pptx(file_path):
        # Load the presentation
        prs = Presentation(file_path)

        # Prepare a list to hold all the text
        all_text = []

        # Iterate through each slide
        for slide in prs.slides:
            # Iterate through each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Append the text of each shape to the list
                    all_text.append(shape.text)

        # Join all the text pieces into a single string
        full_text = "\n".join(all_text)
        return full_text

    

    @staticmethod
    def save_text_to_txt(text, txt_path):
        with open(txt_path, 'w', encoding='utf-8') as file:
            file.write(text)

    @staticmethod
    def split_text(text, max_chunk_size=3800):
        sentences = re.split(r'(?<=[.!?]) +', text)
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > max_chunk_size:
                yield current_chunk
                current_chunk = sentence
            else:
                current_chunk += " " + sentence
        yield current_chunk

    

    @staticmethod
    def create_json_with_Long_summary(json_file_path, summary):
        title_match = re.match(r"([^.]*.)", summary)
        title = title_match.group(0).strip() if title_match else "Summary"
        summary_data = {
            'summary': summary
        }
        with open(json_file_path, 'w') as file:
            json.dump(summary_data, file, indent=4)

    @staticmethod
    def extract_text_from_pdf_plumber(pdf_path, txt_file_path):
        if not os.path.exists(txt_file_path):
            with pdfplumber.open(pdf_path) as pdf:
                with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            txt_file.write(text)
                print(f"Text extracted and saved to {txt_file_path}")     

    @staticmethod
    def extract_images_from_pdf(pdf_path):
        # Open the PDF file
        doc = fitz.open(pdf_path)

        # Extract the PDF name without the extension and create a folder
        pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
        folder_path = f"assets/output_files/Images/{pdf_name}"
        os.makedirs(folder_path, exist_ok=True)

        # Iterate through each page of the PDF
        for page_num in range(len(doc)):
            # Get the page
            page = doc.load_page(page_num)

            # Extract images
            image_list = page.get_images(full=True)

            # Save images
            for image_index, img in enumerate(image_list, start=1):
                # Get the XREF of the image
                xref = img[0]

                # Extract the image bytes
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                # Get the image extension
                image_ext = base_image["ext"]

                # Construct the image filename
                image_filename = f"{pdf_name}_page_{page_num+1}_img_{image_index}.{image_ext}"

                # Save the image
                with open(os.path.join(folder_path, image_filename), "wb") as image_file:
                    image_file.write(image_bytes)

        print(f"Images extracted and saved in folder: {folder_path}")
    @staticmethod
    def extract_images_from_docx(docx_path, output_images_dir):
        doc = Document(docx_path)
        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)
        image_index = 1
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_data = rel.target_part.blob
                image = Image.open(BytesIO(image_data))
                image.save(os.path.join(output_images_dir, f'image_{image_index}.png'))
                image_index += 1
    @staticmethod
    def extract_images_from_pptx(pptx_path):
        # Extract the base name for the pptx file (without extension)
        base_name = DocumentProcessedController.getFileNameFromPathWithOutExtension(pptx_path)
        
        # Set the output directory to the base name of the file
        output_images_dir = "assets/output_files/Images/" + base_name

        # Load the presentation
        prs = Presentation(pptx_path)

        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)

        # Initialize image index
        image_index = 1

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image:
                    # Extract image
                    image_bytes = shape.image.blob
                    image_data = BytesIO(image_bytes)

                    # Open and save the image
                    img = Image.open(image_data)
                    img.save(os.path.join(output_images_dir, f'image_{image_index}.png'), 'PNG')

                    # Increment the image index
                    image_index += 1

        print(f"Images extracted and saved in folder: {output_images_dir}")
    @staticmethod
    def convert_word_to_pdf(docx_path):
        # Determine the output folder based on the docx_path
        output_folder = os.path.dirname(docx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(docx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:writer_pdf_Export",
            "--outdir",
            output_folder,
            docx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
        return output_pdf_path
    @staticmethod
    def convert_ppt_to_pdf(pptx_path):
        # Determine the output folder based on the pptx_path
        output_folder = os.path.dirname(pptx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:impress_pdf_Export",
            "--outdir",
            output_folder,
            pptx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
        return output_pdf_path
    def Document_Processing(self,file):
        if os.path.isfile(file) and file.endswith('.pdf'):
            PDFFile=file
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'assets/output_files/text_files/{filename}.txt'
          
            DocumentProcessedController.extract_text_from_pdf_plumber(PDFFile,text_file_path)
            text =DocumentProcessedController. read_text_file(text_file_path)
            
            
            model = Summarizer()
            text=DocumentProcessedController.clean_text(text)
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}.json") 
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            flashcard=FlashcardsController(text_file_path)
            mcq=QuestionController(PDFFile)
            DocumentProcessedController.extract_images_from_pdf(file)

            #processd_material_id=uuid.uuid4().hex
            #summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            #text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            #documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            #documentprocessed.addProcessedMaterialToFirestore()
        elif os.path.isfile(file) and (file.endswith('.ppt') or file.endswith('.pptx')or file.endswith('.ppsx')):
            file=DocumentProcessedController.convert_ppt_to_pdf(file)
            file=file.replace("\\","/")
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'assets/output_files/text_files/{filename}.txt'
          
            DocumentProcessedController.extract_text_from_pdf_plumber(file,text_file_path)
            text =DocumentProcessedController. read_text_file(text_file_path)
            
            
            model = Summarizer()
            text=DocumentProcessedController.clean_text(text)
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}.json") 
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            flashcard=FlashcardsController(text_file_path)
            mcq=QuestionController(file)
            DocumentProcessedController.extract_images_from_pdf(file)
            # processd_material_id=uuid.uuid4().hex
            # summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            # text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            # documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            # documentprocessed.addProcessedMaterialToFirestore()
        elif os.path.isfile(file) and (file.endswith('.doc') or file.endswith('.docx')):
            file=DocumentProcessedController.convert_word_to_pdf(file)
            file=file.replace("\\","/")
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'assets/output_files/text_files/{filename}.txt'
          
            DocumentProcessedController.extract_text_from_pdf_plumber(file,text_file_path)
            text =DocumentProcessedController. read_text_file(text_file_path)
            
            
            model = Summarizer()
            text=DocumentProcessedController.clean_text(text)
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}.json") 
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            flashcard=FlashcardsController(text_file_path)
            mcq=QuestionController(file)
            DocumentProcessedController.extract_images_from_pdf(file)
            # processd_material_id=uuid.uuid4().hex
            # summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            # text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            # documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            # documentprocessed.addProcessedMaterialToFirestore()
        elif os.path.isfile(file) and (file.endswith('.txt') ):
            txt_path = file
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text=DocumentProcessedController.extract_text_from_word(txt_path)
            text_file = 'f"assets/output_files/text_files/{filename}.txt'
            with open(text_file, 'w') as file:
                file.write(text)
            model = Summarizer()
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json")
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            flashcard=FlashcardsController(text_file)

def main():
    D=DocumentProcessedController("assets/input_files/text-based/test2.pdf")
if __name__ == "__main__":
    main()