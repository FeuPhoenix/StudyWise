import uuid
from app.Studywise.Model import DocumentProcessed
from datetime import datetime
from firebase_admin import firestore
from Constants.Constants import OPENAI_API_KEY, MAX_TOKENS_PER_REQUEST,kUserId,kUserEmail ,kDatejoined ,kFullName 
import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
from PIL import Image
from io import BytesIO
from summarizer import Summarizer
from firebase_admin import credentials, storage
class DocumentProcessedController:
    def __init__(self,material) :
        self.material=material
        self.db = firestore.client()
    @staticmethod
    async def create_processed_document(request_data):
        """
        Creates a new processed document record in Firestore.
        :param request_data: The data from the request to create a new processed document.
        """
        try:
            # Extracting necessary information from the request data
            processed_material_id = request_data.get('processed_material_id')
            material_id = request_data.get('material_id')
            generated_summary_file_path = request_data.get('generated_summary_file_path')
            generated_text_file_path = request_data.get('generated_text_file_path')
            generated_images_file_path = request_data.get('generated_images_file_path')
            generated_document_file_path = request_data.get('generated_document_file_path')
            
            # Creating a new DocumentProcessed instance
            new_document = DocumentProcessed(
                processed_material_id=processed_material_id,
                material_id=material_id,
                generated_summary_file_path=generated_summary_file_path,
                generated_text_file_path=generated_text_file_path,
                generated_images_file_path=generated_images_file_path,
                generated_Document_file_path=generated_document_file_path
            )
            
            # Adding the new document to Firestore
            await new_document.addProcessedMaterialToFirestore()
            
            return {"success": True, "message": "Processed document added successfully."}
        except Exception as e:
            return {"success": False, "message": str(e)}

    @staticmethod
    async def get_processed_document(processed_material_id):
        """
        Retrieves a processed document record from Firestore.
        :param processed_material_id: The ID of the processed document to retrieve.
        """
        try:
            # Logic to retrieve a processed document by ID from Firestore
            # Placeholder for actual Firestore retrieval logic
            pass
            # Assuming the document is successfully retrieved and stored in 'document_data'
            # return {"success": True, "data": document_data}
        except Exception as e:
            return {"success": False, "message": str(e)}

    # Additional methods for updating and deleting processed documents can be added here
    @staticmethod
    def getFileNameFromPathWithOutExtension(input_string):
        last_slash_index = input_string.rfind('/')
        result_string = input_string[last_slash_index + 1:]
        result_string=result_string.replace('.mp4','')
        result_string=result_string.replace('.docx','')
        result_string=result_string.replace('.doc','')
        result_string=result_string.replace('.pptx','')
        result_string=result_string.replace('.ppt','')
        result_string=result_string.replace('.pdf','')
        return result_string
    @staticmethod
    def delete_page_breaks(self,file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        cleaned_content = content.replace("--- Page Break ---", "")
        with open(file_path, 'w') as file:
            file.write(cleaned_content)

    @staticmethod
    def clean_text(text):
        allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
        cleaned_text = re.sub(allowed_pattern, '', text)    
        return cleaned_text

    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r',encoding='utf-8', errors='ignore') as file:
            return file.read()

    @staticmethod
    def extract_text_from_word(file_path):
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        text = DocumentProcessed.clean_text(text)
        return text

    @staticmethod
    def extract_text_from_pptx(prs):
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text

    @staticmethod
    def extract_images_from_pptx(prs, output_images_dir):
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)
        image_index = 1
        for slide_number, slide in enumerate(prs.slides, start=1):
            for shape_number, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    image_data = BytesIO(image_bytes)
                    img = Image.open(image_data)
                    img.save(os.path.join(output_images_dir, f'image_{slide_number}_{shape_number}.png'))
                    image_index += 1

    @staticmethod
    def save_text_to_txt(text, txt_path):
        with open(txt_path, 'w', encoding='utf-8') as file:
            file.write(text)

    @staticmethod
    def split_text(text, max_chunk_size=3800):
        sentences = re.split(r'(?<=[.!?]) +', text)
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > max_chunk_size:
                yield current_chunk
                current_chunk = sentence
            else:
                current_chunk += " " + sentence
        yield current_chunk

    @staticmethod
    def get_Long_summary_from_gpt3(self,text, api_key):
        openai.api_key = api_key
        summaries = []
        for chunk in DocumentProcessedController.split_text(text):
            while True:
                try:
                    response = openai.ChatCompletion.create(
                        model="gpt-3.5-turbo",
                        messages=[
                            {"role": "system", "content": "You are a helpful assistant."},
                            {"role": "user", "content": f"Summarize the following text into concise bullet points, focusing only on essential information: \n{chunk}"}
                        ]
                    )
                    summaries.append(response['choices'][0]['message']['content'].strip())
                    break
                except openai.error.RateLimitError as e:
                    print("Rate limit exceeded, waiting to retry...")
                    time.sleep(20)
        full_summary = ' '.join(summaries)
        return full_summary

    @staticmethod
    def create_json_with_Long_summary(self,json_file_path, summary):
        title_match = re.match(r"([^.]*.)", summary)
        title = title_match.group(0).strip() if title_match else "Summary"
        summary_data = {
            'summary': summary
        }
        with open(json_file_path, 'w') as file:
            json.dump(summary_data, file, indent=4)

    @staticmethod
    def extract_text_from_pdf_plumber(self,pdf_path, txt_file_path):
        if not os.path.exists(txt_file_path):
            with pdfplumber.open(pdf_path) as pdf:
                with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            txt_file.write(text)
                print(f"Text extracted and saved to {txt_file_path}")     
    @staticmethod
    async def upload_to_firebase(local_file, cloud_file):
    # Reference to the storage bucket
        bucket = storage.bucket()

    # Name of the file in the storage bucket
        blob = bucket.blob(cloud_file)

    # Upload the file
        blob.upload_from_filename(local_file)
        
        print(f'{local_file} has been uploaded to {cloud_file}.')
        metadata = blob.metadata
        print(metadata)
        return metadata
    def is_Document(self,file):
        if os.path.isfile(file) and file.endswith('.pdf'):
            PDFFile=file
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'Summaries/transcribed_text_From_{filename}.txt'
            json_file_path = f'Summaries/summary_From_PDF{filename}.json'
            api_key = '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'
            DocumentProcessedController.extract_text_from_pdf_plumber(PDFFile,text_file_path)
            text =DocumentProcessedController. read_text_file(text_file_path)
            text_file = 'f"assets/output_files/text_files/{filename}.txt'
            with open(filename, 'w') as file:
                file.write(text)
            model = Summarizer()
            text=DocumentProcessedController.clean_text(text)
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            processd_material_id=uuid.uuid4().hex
            summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            documentprocessed.addProcessedMaterialToFirestore()
        elif os.path.isfile(file) and (file.endswith('.ppt') or file.endswith('.pptx')):
            pptx_path = file
            filename=DocumentProcessed.getFileNameFromPathWithOutExtension(file)
            output_text_path = f'assets/output_files/{filename}.txt'
            json_file_path = f'Summaries/summary_From_PPT{file}.json'
            prs = Presentation(pptx_path)
            text = DocumentProcessed.extract_text_from_pptx(prs)
            text_file = 'f"assets/output_files/text_files/{filename}.txt'
            with open(filename, 'w') as file:
                file.write(text)
            model = Summarizer()
            text=DocumentProcessed.clean_text(text)
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json") 
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            processd_material_id=uuid.uuid4().hex
            summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            documentprocessed.addProcessedMaterialToFirestore()
        elif os.path.isfile(file) and (file.endswith('.doc') or file.endswith('.docx')):
            docx_path = file
            filename=DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text=DocumentProcessedController.extract_text_from_word(docx_path)
            text_file = 'f"assets/output_files/text_files/{filename}.txt'
            with open(filename, 'w') as file:
                file.write(text)
            model = Summarizer()
            result = model(text)
            summary_data = {
                'long_summary': result
            }
            with open(f"assets/output_files/Summaries/{filename}_summary.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in assets/output_files/Summaries/{filename}_summary.json")
                summary_jsonfile=f"assets/output_files/Summaries/{filename}_summary.json" 
            processd_material_id=uuid.uuid4().hex
            summary_jsonfile_path_on_firebase=DocumentProcessed.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            text_file=DocumentProcessed.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            documentprocessed=DocumentProcessed(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            documentprocessed.addProcessedMaterialToFirestore()