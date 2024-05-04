import subprocess
import uuid
from datetime import datetime
from firebase_admin import firestore, credentials, storage, initialize_app
import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
from summarizer import Summarizer
from backend.Flash_Cards_Controller import FlashcardsController
from backend.DocumentProcessed_Repo import DocumentProcessed
from backend.Question_Controller import QuestionController
import fitz  # PyMuPDF
import comtypes.client
from backend.Constants import OPENAI_API_KEY

from config import socketio

class DocumentProcessedController:
    def __init__(self, file, socketID) :
        self.file = file
        self.socketID = socketID
        self.Document_Processing(file)
    
    # Additional methods for updating and deleting processed documents can be added here
    @staticmethod
    def getFileNameFromPathWithOutExtension(input_string):
        # Extract the base name of the file (i.e., file name with extension, without any directory path)
        base_name_with_extension = os.path.basename(input_string)
        # Split the base name and its extension and return just the name part
        file_name_without_extension, _ = os.path.splitext(base_name_with_extension)
        return file_name_without_extension

    @staticmethod
    def clean_text(text):
        allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
        cleaned_text = re.sub(allowed_pattern, '', text)    
        return cleaned_text

    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r',encoding='utf-8', errors='ignore') as file:
            return file.read()

    @staticmethod
    def extract_text_from_word(file_path):
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        text = DocumentProcessedController.clean_text(text)
        return text

    @staticmethod
    def extract_text_from_pptx(file_path):
        # Load the presentation
        prs = Presentation(file_path)

        # Prepare a list to hold all the text
        all_text = []

        # Iterate through each slide
        for slide in prs.slides:
            # Iterate through each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Append the text of each shape to the list
                    all_text.append(shape.text)

        # Join all the text pieces into a single string
        full_text = "\n".join(all_text)
        return full_text

    

    @staticmethod
    def save_text_to_txt(text, txt_path):
        with open(txt_path, 'w', encoding='utf-8') as file:
            file.write(text)

    @staticmethod
    def split_text(text, max_chunk_size=3800):
        sentences = re.split(r'(?<=[.!?]) +', text)
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > max_chunk_size:
                yield current_chunk
                current_chunk = sentence
            else:
                current_chunk += " " + sentence
        yield current_chunk

    

    @staticmethod
    def create_json_with_Long(json_file_path, summary):
        title_match = re.match(r"([^.]*.)", summary)
        title = title_match.group(0).strip() if title_match else "Summary"
        summary_data = {
            'summary': summary
        }
        with open(json_file_path, 'w') as file:
            json.dump(summary_data, file, indent=4)

    @staticmethod
    def extract_text_from_pdf_plumber(pdf_path, txt_file_path):
        if not os.path.exists(txt_file_path):
            with pdfplumber.open(pdf_path) as pdf:
                os.makedirs(os.path.dirname(txt_file_path), exist_ok=True)
                with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            txt_file.write(text)
                print(f"Text extracted and saved to {txt_file_path}")

    @staticmethod
    def extract_images_from_pdf(pdf_path):
        # Open the PDF file
        doc = fitz.open(pdf_path)

        # Extract the PDF name without the extension and create a folder
        pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
        folder_path = f"mainServerTest/assets/output_files/Images/{pdf_name}"
        os.makedirs(folder_path, exist_ok=True)

        # Iterate through each page of the PDF
        for page_num in range(len(doc)):
            # Get the page
            page = doc.load_page(page_num)

            # Extract images
            image_list = page.get_images(full=True)

            # Save images
            for image_index, img in enumerate(image_list, start=1):
                # Get the XREF of the image
                xref = img[0]

                # Extract the image bytes
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                # Get the image extension
                image_ext = base_image["ext"]

                # Construct the image filename
                image_filename = f"{pdf_name}_page_{page_num+1}_img_{image_index}.{image_ext}"

                # Save the image
                with open(os.path.join(folder_path, image_filename), "wb") as image_file:
                    image_file.write(image_bytes)
        print(f"Images extracted and saved in folder: {folder_path}")

    @staticmethod
    def extract_images_from_docx(docx_path, output_images_dir):
        doc = Document(docx_path)
        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)
        image_index = 1
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_data = rel.target_part.blob
                image = Image.open(BytesIO(image_data))
                image.save(os.path.join(output_images_dir, f'image_{image_index}.png'))
                image_index += 1

    @staticmethod
    def extract_images_from_pptx(pptx_path):
        # Extract the base name for the pptx file (without extension)
        base_name = DocumentProcessedController.getFileNameFromPathWithOutExtension(pptx_path)
        
        # Set the output directory to the base name of the file
        # mainServerTest\assets
        output_images_dir = "mainServerTest/assets/output_files/Images/" + base_name

        # Load the presentation
        prs = Presentation(pptx_path)

        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)

        # Initialize image index
        image_index = 1

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image:
                    # Extract image
                    image_bytes = shape.image.blob
                    image_data = BytesIO(image_bytes)

                    # Open and save the image
                    img = Image.open(image_data)
                    img.save(os.path.join(output_images_dir, f'image_{image_index}.png'), 'PNG')

                    # Increment the image index
                    image_index += 1
        print(f"Images extracted and saved in folder: {output_images_dir}")

    @staticmethod
    def convert_word_to_pdf(docx_path):
        # Determine the output folder based on the docx_path
        output_folder = os.path.dirname(docx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(docx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:writer_pdf_Export",
            "--outdir",
            output_folder,
            docx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
        return output_pdf_path
    
    @staticmethod
    def convert_ppt_to_pdf(pptx_path):
        # Determine the output folder based on the pptx_path
        output_folder = os.path.dirname(pptx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:impress_pdf_Export",
            "--outdir",
            output_folder,
            pptx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
        return output_pdf_path
    
    @staticmethod
    def get_Long_summary(text):
        openai.api_key = OPENAI_API_KEY
        summaries = []

        for chunk in DocumentProcessedController.split_text(text):
            while True:  # Keep trying until successful
                try:
                    response = openai.ChatCompletion.create(
                        model="gpt-3.5-turbo",
                        messages=[
                            {"role": "system", "content": "You are a helpful assistant."},
                            {"role": "user", "content": f"Summarize the following text into concise bullet points, focusing only on essential information: \n{chunk}"}
                        ]
                    )
                    summaries.append(response['choices'][0]['message']['content'].strip())
                    break  # Exit the loop if successful
                except openai.error.RateLimitError as e:
                    print("Rate limit exceeded, waiting to retry...")
                    time.sleep(20)  # Wait for 20 seconds before retrying

        full_summary = ' '.join(summaries)
        return full_summary

    def Document_Processing(self, file):
        if os.path.isfile(file) and file.endswith('.pdf'):

            socketio.emit('update', {'message': 'PDF accepted.'}, to=self.socketID)

            PDFFile = file
            filename = DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            file = f'mainServerTest/assets/input_files/text_based/{filename}.pdf'
            text_file_path = f'mainServerTest/assets/output_files/text_files/{filename}.txt'

            socketio.emit('update', {'message': 'Extracting text..'}, to=self.socketID)
            DocumentProcessedController.extract_text_from_pdf_plumber(PDFFile, text_file_path)
            text = DocumentProcessedController.read_text_file(text_file_path)

            socketio.emit('update', {'message': 'Formatting text..'}, to=self.socketID)
            text = DocumentProcessedController.clean_text(text)

            socketio.emit('update', {'message': 'Generating Summary..'}, to=self.socketID)
            result = DocumentProcessedController.get_Long_summary(text)

            summary_data = {
                'long': result
            }
            print(f"filename is: {filename}")
            with open(f"mainServerTest/assets/output_files/summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in mainServerTest/assets/output_files/summaries/{filename}.json") 
                summary_jsonfile = f"mainServerTest/assets/output_files/summaries/{filename}.json" 
            
            socketio.emit('update', {'message': 'Generating Flashcards..'}, to=self.socketID)
            flashcard = FlashcardsController(text_file_path)

            socketio.emit('update', {'message': 'Generating MCQs..'}, to=self.socketID)
            mcq = QuestionController(PDFFile)

            socketio.emit('update', {'message': 'Extracting images..'}, to=self.socketID)
            DocumentProcessedController.extract_images_from_pdf(file)

            #processd_material_id=uuid.uuid4().hex
            #summary_jsonfile_path_on_firebase=DocumentProcessedController.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            #text_file=DocumentProcessedController.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            #DocumentProcessedController=DocumentProcessedController(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            #DocumentProcessedController.addProcessedMaterialToFirestore()

        elif os.path.isfile(file) and (file.endswith('.ppt') or file.endswith('.pptx') or file.endswith('.ppsx')):
            socketio.emit('update', {'message': 'Powerpoint accepted. Converting to PDF..'}, to=self.socketID)
            file = DocumentProcessedController.convert_ppt_to_pdf(file)

            file = file.replace("\\","/")
            filename = DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'mainServerTest/assets/output_files/text_files/{filename}.txt'

            socketio.emit('update', {'message': 'Extracting text..'}, to=self.socketID)
            DocumentProcessedController.extract_text_from_pdf_plumber(file, text_file_path)
            text = DocumentProcessedController.read_text_file(text_file_path)

            socketio.emit('update', {'message': 'Formatting text..'}, to=self.socketID)
            text = DocumentProcessedController.clean_text(text)

            socketio.emit('update', {'message': 'Generating Summary..'}, to=self.socketID)
            result = DocumentProcessedController.get_Long_summary(text)

            summary_data = {
                'long': result
            }
            with open(f"mainServerTest/assets/output_files/summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in mainServerTest/assets/output_files/summaries/{filename}.json") 
                summary_jsonfile=f"mainServerTest/assets/output_files/summaries/{filename}.json"
            
            socketio.emit('update', {'message': 'Generating Flashcards..'}, to=self.socketID)
            flashcard = FlashcardsController(text_file_path)

            socketio.emit('update', {'message': 'Generating MCQs..'}, to=self.socketID)
            mcq = QuestionController(file)

            socketio.emit('update', {'message': 'Extracting images..'}, to=self.socketID)
            DocumentProcessedController.extract_images_from_pdf(file)

            # processd_material_id=uuid.uuid4().hex
            # summary_jsonfile_path_on_firebase=DocumentProcessedController.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            # text_file=DocumentProcessedController.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            # DocumentProcessedController=DocumentProcessedController(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            # DocumentProcessedController.addProcessedMaterialToFirestore()

        elif os.path.isfile(file) and (file.endswith('.doc') or file.endswith('.docx')):
            socketio.emit('update', {'message': 'Accepted Word document. Converting to PDF'}, to=self.socketID)
            file = DocumentProcessedController.convert_word_to_pdf(file)

            file = file.replace("\\","/")
            filename = DocumentProcessedController.getFileNameFromPathWithOutExtension(file)
            text_file_path = f'mainServerTest/assets/output_files/text_files/{filename}.txt'
          
            socketio.emit('update', {'message': 'Extracting text..'}, to=self.socketID)
            DocumentProcessedController.extract_text_from_pdf_plumber(file, text_file_path)
            text = DocumentProcessedController.read_text_file(text_file_path)

            socketio.emit('update', {'message': 'Formatting text..'}, to=self.socketID)
            text = DocumentProcessedController.clean_text(text)

            socketio.emit('update', {'message': 'Generating Summary..'}, to=self.socketID)
            result = DocumentProcessedController.get_Long_summary(text)

            summary_data = {
                'long': result
            }
            with open(f"mainServerTest/assets/output_files/summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in mainServerTest/assets/output_files/summaries/{filename}.json") 
                summary_jsonfile = f"mainServerTest/assets/output_files/summaries/{filename}.json" 

            socketio.emit('update', {'message': 'Generating Flashcards..'}, to=self.socketID)
            flashcard = FlashcardsController(text_file_path)

            socketio.emit('update', {'message': 'Generating MCQs..'}, to=self.socketID)
            mcq = QuestionController(file)

            socketio.emit('update', {'message': 'Exctracting images..'}, to=self.socketID)
            DocumentProcessedController.extract_images_from_pdf(file)

            # processd_material_id=uuid.uuid4().hex
            # summary_jsonfile_path_on_firebase=DocumentProcessedController.upload_to_firebase(summary_jsonfile,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(summary_jsonfile)}.json')
            # text_file=DocumentProcessedController.upload_to_firebase(text_file,f'{kUserId}/{DocumentProcessedController.getFileNameFromPathWithOutExtension(text_file)}.txt')
            # DocumentProcessedController=DocumentProcessedController(processd_material_id,self.material,summary_jsonfile_path_on_firebase,text_file) 
            # DocumentProcessedController.addProcessedMaterialToFirestore()

        elif os.path.isfile(file) and (file.endswith('.txt') ):
            socketio.emit('update', {'message': 'Accepted text file (.txt)'}, to=self.socketID)
            txt_path = file
            filename = DocumentProcessedController.getFileNameFromPathWithOutExtension(file)

            socketio.emit('update', {'message': 'Extracting text..'}, to=self.socketID)
            text = DocumentProcessedController.extract_text_from_word(txt_path)

            text_file = f'mainServerTest/assets/output_files/text_files/{filename}.txt'
            with open(text_file, 'w') as file:
                file.write(text)

            socketio.emit('update', {'message': 'Generating Summary..'}, to=self.socketID)
            result = DocumentProcessedController.get_Long_summary(text)

            summary_data = {
                'long': result
            }
            with open(f"mainServerTest/assets/output_files/summaries/{filename}.json", 'w') as json_file:
                json.dump(summary_data, json_file, indent=4)
                print(f"Long summary has been successfully saved in mainServerTest/assets/output_files/summaries/{filename}.json")
                summary_jsonfile = f"mainServerTest/assets/output_files/summaries/{filename}.json" 
            
            socketio.emit('update', {'message': 'Generating Flashcards..'}, to=self.socketID)
            flashcard = FlashcardsController(text_file)

def main(file, socketID):
    print("Current Working Directory: \n", os.getcwd())
    D = DocumentProcessedController(file, socketID)
if __name__ == "__main__":
    main()