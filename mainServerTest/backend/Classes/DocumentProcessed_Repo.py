import asyncio
from datetime import datetime, timedelta
from firebase_admin import initialize_app, storage
import json
import time
import openai
import re
import pdfplumber
import os
from pptx import Presentation
from io import BytesIO
from PIL import Image
from docx import Document
import subprocess
import uuid
import fitz  # PyMuPDF
import comtypes.client
from dotenv import load_dotenv
from langdetect import detect

# from backend.Classes.Constants import OPENAI_API_KEY 
# from backend.Classes.FirestoreDB import FirestoreDB
# from backend.Classes.Flash_Cards_Repo import Flash_Cards
# from backend.Classes.Questions_Repo import Questions_Repo
from Constants import OPENAI_API_KEY 
from FirestoreDB import FirestoreDB
from Flash_Cards_Repo import Flash_Cards
from Questions_Repo import Questions_Repo

from dotenv import load_dotenv
load_dotenv()

class DocumentProcessed:
    #material_id=uuid.uuid4().hex#done
    # file_name=""#done
    # _file_path=""#done
    # file_type=""#done
    # generated_summary_file_path=""
    # generated_text_file_path=""
    # generated_text_file_path=""
    # generated_images_folder_path=""
    # userid=""
  
    def __init__(self, file,UserID):
        self.file=file
        self.userid=UserID
        self.material_id=uuid.uuid4().hex
        self.db = FirestoreDB.get_instance()
        
    @staticmethod
    def detect_language(text):
        return detect(text)    
    def addProcessedMaterialToFirestore(self):
        db_instance = FirestoreDB.get_instance()
        firestore_instance = db_instance.get_firestore_instance()
        
       
        file_path_location,generated_text_file_path,generated_summary_file_path=DocumentProcessed.upload_material_to_storage(self.userid,self.file_name , self._file_path,self.generated_text_file_path,self.generated_summary_file_path)

        #document('13ffe4704e2d423ea7751cb88d599db7') the number will be replaced with the user id
        #document(rmk3SGTciwNRdo9pT4CO) this will be replaced with the material id


        try:
            doc_ref=firestore_instance.collection('Users').document(self.userid).collection('DocumentMaterial').document(self.material_id).set({
                "file_name": self.file_name,
                "_file_path": file_path_location,
                "file_type": self.file_type,
                "material_id": self.material_id,
                "Material":self._file_path,
                "generated_summary_file_path": generated_summary_file_path,
                "generated_text_file_path":generated_text_file_path ,
                #"generated_images_folder_path": DocumentProcessed.upload_file_to_storage(self._file_path, f"Users/Uploaded Materials/{self.file_name}",self.generated_images_folder_path),

            })
            print("Successfully added Document to firestore")
        except Exception as e:
            print(e)

    @classmethod
    def fromJson(cls, data):
        return cls(
            material_id=data["material_id"],
            user_ID=data["user_ID"],
            file_name=data["file_name"],
            _file_path=data["_file_path"],
            file_type=data["file_type"],
            generated_summary_file_path=data["generated_summary_file_path"],
            generated_text_file_path=data["generated_text_file_path"],
            generated_images_folder_path=data["generated_images_folder_path"],
            generated_flashcards_file=data["generated_flashcards_file"],
            generated_MCQ_file_E=data["generated_MCQ_file_E"],
            generated_MCQ_file_M=data["generated_MCQ_file_M"],
            generated_MCQ_file_H=data["generated_MCQ_file_H"],
            
        )
    def get_Document_from_firestore(self):
            db_instance = FirestoreDB.get_instance()
            firestore_instance = db_instance.get_firestore_instance()
            try:
            # Reference to the document
                doc_ref = firestore_instance.collection('Users').document('13ffe4704e2d423ea7751cb88d599db7').collection('DocumentMaterial').document("rmk3SGTciwNRdo9pT4CO")
                
                # Get the document snapshot
                doc = doc_ref.get()
                
                # Check if the document exists
                if doc.exists:
                    # Get the data from the document
                    data = doc.to_dict()
                    print("Document Data:", data)  # Print the data

                    return data
                else:
                    print("Document does not exist")
                    return None
            except Exception as e:
                print("Error:", e)
                return None
    @staticmethod
    def upload_material_to_storage(user_id, material_name, material_file_path, text_file_path, summary_file_path):
        db_instance = FirestoreDB.get_instance()
        storage_instance = db_instance.get_storage_instance()

        # Constructing the full path for the folder using the material name
        folder_path = f"user/{user_id}/Uploaded Materials/{material_name}/"

        # Creating a folder with the material name as the folder name
        folder_blob = storage_instance.blob(folder_path)
        folder_blob.upload_from_string('')  # Upload an empty string to create the folder

        # Upload the material file inside the folder
        material_blob_path = folder_path + "material"
        material_blob = storage_instance.blob(material_blob_path)
        material_blob.upload_from_filename(material_file_path,timeout=600)

        # Upload the text file inside the folder
        text_blob_path = folder_path + "text.txt"
        text_blob = storage_instance.blob(text_blob_path)
        text_blob.upload_from_filename(text_file_path,timeout=600)

        # Upload the summary file inside the folder
        summary_blob_path = folder_path + "summary.txt"
        summary_blob = storage_instance.blob(summary_blob_path)
        summary_blob.upload_from_filename(summary_file_path,timeout=600)
        
        expiration = datetime.now() + timedelta(days=36500)
        # download_urls = {
        #     "material_url": material_blob.generate_signed_url(expiration=expiration),
        #     "text_url": text_blob.generate_signed_url(expiration=expiration),
        #     "summary_url": summary_blob.generate_signed_url(expiration=expiration)
        # }

        print("Successfully uploaded material to Storage")
        return material_blob.generate_signed_url(expiration=expiration),text_blob.generate_signed_url(expiration=expiration),summary_blob.generate_signed_url(expiration=expiration)

    def toJson(self):
        data = {
            "material_id": self.material_id,
            "user_ID": self.userid,
            "user_ID": self.userid,
            "file_name": self.file_name,
            "_file_path": self._file_path,
            "file_type": self.file_type,
            "generated_summary_file_path": self.generated_summary_file_path,
            "generated_text_file_path": self.generated_text_file_path,
            "generated_images_folder_path": self.generated_images_folder_path,

        }
        return data
    async def upload_to_firebase(self,local_file, cloud_file):
    # Reference to the storage bucket
        bucket = storage.bucket()

    # Name of the file in the storage bucket
        blob = bucket.blob(cloud_file)

    # Upload the file
        await blob.upload_from_filename(local_file)
        
        print(f'{local_file} has been uploaded to {cloud_file}.')
        metadata = blob.metadata
        print(metadata)
        return metadata


    @staticmethod
    def getFileNameFromPathWithOutExtension(input_string):
        last_slash_index = input_string.rfind('\\')
        result_string = input_string[last_slash_index + 1:]
        result_string=result_string.replace('.mp4','')
        result_string=result_string.replace('.docx','')
        result_string=result_string.replace('.doc','')
        result_string=result_string.replace('.pptx','')
        result_string=result_string.replace('.ppt','')
        result_string=result_string.replace('.pdf','')
        return result_string

    @staticmethod
    def clean_text(text):
        allowed_pattern = r"[^a-zA-Z\s,.;:'\"!?()-]"
        cleaned_text = re.sub(allowed_pattern, '', text)    
        return cleaned_text

    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r',encoding='utf-8', errors='ignore') as file:
            return file.read()

    @staticmethod
    def extract_text_from_word(file_path):
        doc = Document(file_path)
        text = ''
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        text = DocumentProcessed.clean_text(text)
        return text

    @staticmethod
    def extract_text_from_pptx(file_path):
        # Load the presentation
        prs = Presentation(file_path)

        # Prepare a list to hold all the text
        all_text = []

        # Iterate through each slide
        for slide in prs.slides:
            # Iterate through each shape in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Append the text of each shape to the list
                    all_text.append(shape.text)

        # Join all the text pieces into a single string
        full_text = "\n".join(all_text)
        return full_text

    

    @staticmethod
    def save_text_to_txt(text, txt_path):
        with open(txt_path, 'w', encoding='utf-8') as file:
            file.write(text)

    @staticmethod
    def split_text(text, max_chunk_size=3800):
        sentences = re.split(r'(?<=[.!?]) +', text)
        current_chunk = ""
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > max_chunk_size:
                yield current_chunk
                current_chunk = sentence
            else:
                current_chunk += " " + sentence
        yield current_chunk

    

    @staticmethod
    def create_json_with_Long_summary(json_file_path, summary):
        title_match = re.match(r"([^.]*.)", summary)
        title = title_match.group(0).strip() if title_match else "Summary"
        summary_data = {
            'summary': summary
        }
        with open(json_file_path, 'w') as file:
            json.dump(summary_data, file, indent=4)

    @staticmethod
    def extract_text_from_pdf_plumber(pdf_path, txt_file_path):
        if not os.path.exists(txt_file_path):
            with pdfplumber.open(pdf_path) as pdf:
                with open(txt_file_path, 'w', encoding='utf-8') as txt_file:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            txt_file.write(text)
                print(f"Text extracted and saved to {txt_file_path}")     

    @staticmethod
    def extract_images_from_pdf(pdf_path):
        # Open the PDF file
        doc = fitz.open(pdf_path)

        # Extract the PDF name without the extension and create a folder
        pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
        folder_path = f"mainServerTest/assets/output_files/Images/{pdf_name}"
        os.makedirs(folder_path, exist_ok=True)

        # Iterate through each page of the PDF
        for page_num in range(len(doc)):
            # Get the page
            page = doc.load_page(page_num)

            # Extract images
            image_list = page.get_images(full=True)

            # Save images
            for image_index, img in enumerate(image_list, start=1):
                # Get the XREF of the image
                xref = img[0]

                # Extract the image bytes
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]

                # Get the image extension
                image_ext = base_image["ext"]

                # Construct the image filename
                image_filename = f"{pdf_name}_page_{page_num+1}_img_{image_index}.{image_ext}"

                # Save the image
                with open(os.path.join(folder_path, image_filename), "wb") as image_file:
                    image_file.write(image_bytes)

        print(f"Images extracted and saved in folder: {folder_path}")
    @staticmethod
    def extract_images_from_docx(docx_path, output_images_dir):
        doc = Document(docx_path)
        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)
        image_index = 1
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_data = rel.target_part.blob
                image = Image.open(BytesIO(image_data))
                image.save(os.path.join(output_images_dir, f'image_{image_index}.png'))
                image_index += 1
    @staticmethod
    def extract_images_from_pptx(pptx_path):
        # Extract the base name for the pptx file (without extension)
        base_name = DocumentProcessed.getFileNameFromPathWithOutExtension(pptx_path)
        
        # Set the output directory to the base name of the file
        output_images_dir = "mainServerTest/assets/output_files/Images/" + base_name

        # Load the presentation
        prs = Presentation(pptx_path)

        # Create the output images directory if it doesn't exist
        if not os.path.exists(output_images_dir):
            os.makedirs(output_images_dir)

        # Initialize image index
        image_index = 1

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "image") and shape.image:
                    # Extract image
                    image_bytes = shape.image.blob
                    image_data = BytesIO(image_bytes)

                    # Open and save the image
                    img = Image.open(image_data)
                    img.save(os.path.join(output_images_dir, f'image_{image_index}.png'), 'PNG')

                    # Increment the image index
                    image_index += 1

        print(f"Images extracted and saved in folder: {output_images_dir}")
    @staticmethod
    
    def check_value_exists_in_DocumentMaterial(userid, attribute_value):
        db_instance = FirestoreDB.get_instance()  # Assuming FirestoreDB is a class or method to access Firestore
        firestore_instance = db_instance.get_firestore_instance()
        try:
            # Reference to the collection
            doc_material_ref = firestore_instance.collection("Users").document(userid).collection("DocumentMaterial")

            # Query documents where the attribute name matches the value
            query = doc_material_ref.where("file_name", "==", attribute_value).stream()

            # Iterate over the query results
            for doc in query:
                # If any document matches the query, return False
                print("The Document already exists")
                return False

            # If no matching document is found, return True
            return True
        except Exception as e:
            print("Error:", e)
            return None
    @staticmethod
    def convert_word_to_pdf(docx_path):
        # Determine the output folder based on the docx_path
        output_folder = os.path.dirname(docx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(docx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:writer_pdf_Export",
            "--outdir",
            output_folder,
            docx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
        return output_pdf_path
    @staticmethod
    def convert_txt_to_pdf(txt_path):
        # Determine the output folder based on the txt_path
        output_folder = os.path.dirname(txt_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(txt_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_folder,
            txt_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

    # Return the path of the generated PDF
        return output_pdf_path
    @staticmethod
    def convert_ppt_to_pdf(pptx_path):
        # Determine the output folder based on the pptx_path
        output_folder = os.path.dirname(pptx_path)

        # Construct the base name for the output PDF file
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        output_pdf_path = os.path.join(output_folder, f"{base_name}.pdf")

        # Ensure the output folder exists
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        # Construct the command
        command = [
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "--headless",
            "--convert-to",
            "pdf:impress_pdf_Export",
            "--outdir",
            output_folder,
            pptx_path
        ]

        # Execute the command
        subprocess.run(command, shell=True)

        # Return the path of the generated PDF
     
        return output_pdf_path
    
    @staticmethod
    def get_file_extension(file_name):
        # Split the file name by the dot (.)
        parts = file_name.split(".")

        # Check if there is an extension available
        if len(parts) > 1:
            # Return the last part as the extension
            return parts[-1]
        else:
            # No extension found
            return ""
    
    @staticmethod
    def split_text(text, max_chunk_size=3800):  # Reduced max_chunk_size for safety
    
        sentences = re.split(r'(?<=[.!?]) +', text)
        current_chunk = ""
        for sentence in sentences:
            # Check if adding the next sentence would exceed the max chunk size EZ
            if len(current_chunk) + len(sentence) > max_chunk_size:
                yield current_chunk
                current_chunk = sentence  # Start a new chunk EZ
            else:
                current_chunk += " " + sentence
        yield current_chunk  # Yield the last chunk
    
    @staticmethod
    def get_Long_summary(text):
        openai.api_key = OPENAI_API_KEY
        summaries = []

        for chunk in DocumentProcessed.split_text(text):
            while True:  # Keep trying until successful
                try:
                    response = openai.ChatCompletion.create(
                        model="gpt-3.5-turbo",
                        messages=[
                            {"role": "system", "content": "You are a helpful assistant."},
                            {"role": "user", "content": f"Summarize the following text into concise bullet points, focusing only on essential information: \n{chunk}"}
                        ]
                    )
                    summaries.append(response['choices'][0]['message']['content'].strip())
                    break  # Exit the loop if successful
                except openai.error.RateLimitError as e:
                    print("Rate limit exceeded, waiting to retry...")
                    time.sleep(20)  # Wait for 20 seconds before retrying

        full_summary = ' '.join(summaries)
        return full_summary
    
    def get_long_summary_and_write_to_json(self,text, filename):
        result = DocumentProcessed.get_Long_summary(text)
        summary_data = {'long_summary': result}
        with open(f"mainServerTest/assets/output_files/summaries/{filename}.json", 'w') as json_file:
            json.dump(summary_data, json_file, indent=4)
            print(f"Long summary has been successfully saved in mainServerTest/assets/output_files/summaries/{filename}.json") 
        self.generated_summary_file_path = f"mainServerTest/assets/output_files/summaries/{filename}.json" 
    
    def Document_Processing(self):
        self.file_name = DocumentProcessed.getFileNameFromPathWithOutExtension(self.file)

        print("File name without path: ", self.file_name)

        if DocumentProcessed.check_value_exists_in_DocumentMaterial(self.userid, self.file_name):
            if os.path.isfile(self.file):
                if self.file.endswith('.pdf'):
                    self.file_type = DocumentProcessed.get_file_extension(self.file)
                    text_file_path = f'mainServerTest/assets/output_files/text_files/{self.file_name}.txt'
                    self.generated_text_file_path = text_file_path
                    DocumentProcessed.extract_text_from_pdf_plumber(self.file, text_file_path)
                elif self.file.endswith(('.ppt', '.pptx', '.ppsx')):
                    self.file_type = DocumentProcessed.get_file_extension(self.file)
                    self._file_path = DocumentProcessed.convert_ppt_to_pdf(self.file)
                    self.generated_text_file_path = f'mainServerTest/assets/output_files/text_files/{self.file_name}.txt'
                    DocumentProcessed.extract_text_from_pdf_plumber(self._file_path, self.generated_text_file_path)
                elif self.file.endswith(('.doc', '.docx')):
                    self.file_type = DocumentProcessed.get_file_extension(self.file)
                    self._file_path = DocumentProcessed.convert_word_to_pdf(self.file)
                    self.generated_text_file_path = f'mainServerTest/assets/output_files/text_files/{self.file_name}.txt'
                    DocumentProcessed.extract_text_from_pdf_plumber(self._file_path, self.generated_text_file_path)
                elif self.file.endswith('.txt'):
                    self.file_type = DocumentProcessed.get_file_extension(self.file)
                    self._file_path = DocumentProcessed.convert_txt_to_pdf(self.file)
                    self.generated_text_file_path = f'mainServerTest/assets/output_files/text_files/{self.file_name}.txt'
                    text = DocumentProcessed.extract_text_from_word(self.file)
                    with open(self.generated_text_file_path, 'w') as file:
                        file.write(text)
                else:
                    print("Please provide a valid file path")
                text = DocumentProcessed.read_text_file(self.generated_text_file_path)
                text = DocumentProcessed.clean_text(text)
                result = DocumentProcessed.get_Long_summary(text)
                self.get_long_summary_and_write_to_json(result, self.file_name)
                self.addProcessedMaterialToFirestore()
                flashcard = Flash_Cards(self._file_path, self.userid, self.material_id)
                mcq = Questions_Repo(self._file_path, self.userid, self.material_id, None)
                
            else:
                print("Please provide a valid file path")
        else:
            print("The Document already exists")


# Testing
def main():
    D = DocumentProcessed("C:/Users/Abdelrahman/Downloads/محاضرة د.محمود البحيري ٢ (1).pdf", "0GKTloo0geWML96tvd9g27C99543")
    a, b = D.Document_Processing()
    print('Material Filename and ID: \n', a, b)


if __name__ == "__main__":
    main()
