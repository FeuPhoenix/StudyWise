import os
from pptx import Presentation
from io import BytesIO
from PIL import Image

def extract_text_from_pptx(prs):
    text = ''

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'

    return text

def extract_images_from_pptx(prs, output_images_dir):
    # Create the output images directory if it doesn't exist
    if not os.path.exists(output_images_dir):
        os.makedirs(output_images_dir)

    image_index = 1
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_data = BytesIO(image_bytes)
                img = Image.open(image_data)
                img.save(os.path.join(output_images_dir, f'image_{slide_number}_{shape_number}.png'))
                image_index += 1

def save_text_to_txt(text, txt_path):
    # Save the extracted text to a text file
    with open(txt_path, 'w', encoding='utf-8') as file:
        file.write(text)

def main():
    # Replace 'your_presentation.pptx' with the path to your PowerPoint presentation
    pptx_path = 'D:/AA/Sw/testing/test2.pptx'

    # Replace 'output_text.txt' with the desired path for the output text file
    output_text_path = 'D:/AA/Sw/testing/output_text.txt'

    # Replace 'output_images' with the desired directory for the output images
    output_images_dir = 'D:/AA/Sw/testing/output_images2'

    # Load the PowerPoint presentation
    prs = Presentation(pptx_path)

    # Extract text from the PowerPoint presentation
    text = extract_text_from_pptx(prs)

    # Extract images from the PowerPoint presentation
    extract_images_from_pptx(prs, output_images_dir)

    # Save the extracted text to a text file
    save_text_to_txt(text, output_text_path)

if __name__ == "__main__":
    main()
